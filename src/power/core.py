"""Core presentation builder for Power."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from power.slides import SlideBuilder

logger = logging.getLogger(__name__)


class PowerPresentation:
    """Main presentation builder class.

    Example usage:
        >>> ppt = PowerPresentation()
        >>> ppt.add_title_slide("My Presentation", "Subtitle")
        >>> ppt.add_content_slide("Topic", ["Point 1", "Point 2"])
        >>> ppt.save("output.pptx")

    With template:
        >>> ppt = PowerPresentation(template="Galaxy.pptx")
        >>> ppt.add_title_slide("Using Galaxy Theme")
        >>> ppt.save("themed.pptx")
    """

    def __init__(self, template: str | Path | None = None):
        """Initialize a new presentation.

        Args:
            template: Optional path to a .pptx template file.
        """
        self.template_path = Path(template) if template else None
        self._prs: Presentation | None = None
        self._load_presentation()

    def _load_presentation(self) -> None:
        """Load presentation from template or create blank."""
        if self.template_path:
            if not self.template_path.exists():
                raise FileNotFoundError(f"Template not found: {self.template_path}")
            try:
                self._prs = Presentation(str(self.template_path))
                logger.info(f"Loaded template: {self.template_path}")
            except PackageNotFoundError as e:
                raise ValueError(f"Invalid PowerPoint file: {self.template_path}") from e
        else:
            self._prs = Presentation()
            logger.info("Created blank presentation")

    @property
    def presentation(self) -> Presentation:
        """Access the underlying python-pptx Presentation object."""
        if self._prs is None:
            raise RuntimeError("Presentation not initialized")
        return self._prs

    @property
    def slide_count(self) -> int:
        """Return the number of slides in the presentation."""
        return len(self.presentation.slides)

    @property
    def layout_count(self) -> int:
        """Return the number of available slide layouts."""
        return len(self.presentation.slide_layouts)

    def get_layout_names(self) -> list[str]:
        """Get names of all available slide layouts."""
        return [layout.name for layout in self.presentation.slide_layouts]

    def get_layout(self, index_or_name: int | str):
        """Get a slide layout by index or name.

        Args:
            index_or_name: Layout index (int) or name (str).

        Returns:
            SlideLayout object.
        """
        if isinstance(index_or_name, int):
            if 0 <= index_or_name < self.layout_count:
                return self.presentation.slide_layouts[index_or_name]
            raise IndexError(f"Layout index {index_or_name} out of range (0-{self.layout_count-1})")

        # Search by name
        for layout in self.presentation.slide_layouts:
            if layout.name.lower() == index_or_name.lower():
                return layout
        raise ValueError(f"Layout '{index_or_name}' not found")

    def find_content_layout(self) -> int:
        """Find the best layout for content slides (title + content).

        Searches for layouts named 'Title and Content' or similar.
        Falls back to layout 1 for blank presentations.

        Returns:
            Layout index for content slides.
        """
        # Priority search for content layouts by name
        content_layout_names = [
            "title and content",
            "title & content",
            "title, content",
        ]

        for i, layout in enumerate(self.presentation.slide_layouts):
            name_lower = layout.name.lower()
            # Exact match first
            if name_lower in content_layout_names:
                return i
            # Partial match: "Title and Content" but not "Title and Content 3" etc
            if "title and content" in name_lower and name_lower.endswith("content"):
                return i

        # Fallback: look for layout with TITLE and OBJECT/BODY placeholders
        for i, layout in enumerate(self.presentation.slide_layouts):
            has_title = False
            has_content = False
            for ph in layout.placeholders:
                ph_type = str(ph.placeholder_format.type)
                if "TITLE" in ph_type and "CENTER" not in ph_type:
                    has_title = True
                if "OBJECT" in ph_type or "BODY" in ph_type:
                    has_content = True
            if has_title and has_content:
                return i

        # Final fallback: layout 1 (standard for blank presentations)
        return 1

    def add_slide(self, layout: int | str = 6) -> SlideBuilder:
        """Add a new slide and return a SlideBuilder for it.

        Args:
            layout: Layout index or name. Default is 6 (Blank).

        Returns:
            SlideBuilder for the new slide.
        """
        slide_layout = self.get_layout(layout)
        slide = self.presentation.slides.add_slide(slide_layout)
        return SlideBuilder(slide)

    def add_title_slide(
        self, title: str, subtitle: str = "", layout: int | str = 0
    ) -> SlideBuilder:
        """Add a title slide.

        Args:
            title: Main title text.
            subtitle: Optional subtitle text.
            layout: Layout to use (default: 0 = Title Slide).

        Returns:
            SlideBuilder for the new slide.
        """
        builder = self.add_slide(layout)
        builder.set_title(title)
        if subtitle:
            builder.set_subtitle(subtitle)
        return builder

    def add_section_slide(self, title: str, subtitle: str = "") -> SlideBuilder:
        """Add a section header slide.

        Args:
            title: Section title.
            subtitle: Optional subtitle.

        Returns:
            SlideBuilder for the new slide.
        """
        return self.add_title_slide(title, subtitle, layout=2)

    def add_content_slide(
        self, title: str, bullets: list[str] | None = None, layout: int | str | None = None
    ) -> SlideBuilder:
        """Add a content slide with optional bullet points.

        Args:
            title: Slide title.
            bullets: Optional list of bullet points.
            layout: Layout to use. If None, automatically finds "Title and Content" layout.

        Returns:
            SlideBuilder for the new slide.
        """
        if layout is None:
            layout = self.find_content_layout()
        builder = self.add_slide(layout)
        builder.set_title(title)
        if bullets:
            builder.add_bullets(bullets)
        return builder

    def add_blank_slide(self) -> SlideBuilder:
        """Add a blank slide for custom content.

        Returns:
            SlideBuilder for the new slide.
        """
        return self.add_slide(6)

    def remove_slide(self, index: int) -> None:
        """Remove a slide by index.

        Args:
            index: Index of slide to remove (0-based).
        """
        if not 0 <= index < self.slide_count:
            raise IndexError(f"Slide index {index} out of range")

        slide_id = self.presentation.slides._sldIdLst[index].rId
        self.presentation.part.drop_rel(slide_id)
        del self.presentation.slides._sldIdLst[index]

    def clear_slides(self) -> None:
        """Remove all slides from the presentation."""
        while self.slide_count > 0:
            self.remove_slide(0)

    def replace_placeholders(self, replacements: dict[str, str]) -> None:
        """Replace {{PLACEHOLDER}} text throughout presentation.

        Args:
            replacements: Dict mapping placeholder names to values.
                         e.g., {"{{NAME}}": "John", "{{DATE}}": "2024-01-01"}
        """
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for placeholder, value in replacements.items():
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, str(value))

    def get_template_info(self) -> dict[str, Any]:
        """Get information about the presentation/template.

        Returns:
            Dict with template information.
        """
        info = {
            "template": str(self.template_path) if self.template_path else None,
            "slide_count": self.slide_count,
            "layout_count": self.layout_count,
            "layouts": [],
            "slide_width": self.presentation.slide_width,
            "slide_height": self.presentation.slide_height,
        }

        for i, layout in enumerate(self.presentation.slide_layouts):
            layout_info = {
                "index": i,
                "name": layout.name,
                "placeholders": [],
            }
            for ph in layout.placeholders:
                layout_info["placeholders"].append(
                    {
                        "idx": ph.placeholder_format.idx,
                        "type": str(ph.placeholder_format.type),
                        "name": ph.name,
                    }
                )
            info["layouts"].append(layout_info)

        return info

    def save(self, output_path: str | Path) -> Path:
        """Save the presentation to a file.

        Args:
            output_path: Path for the output file.

        Returns:
            Path to the saved file.
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        self.presentation.save(str(output_path))
        logger.info(f"Saved presentation to {output_path}")
        return output_path

    def __enter__(self) -> PowerPresentation:
        """Context manager entry."""
        return self

    def __exit__(self, exc_type, exc_val, exc_tb) -> None:
        """Context manager exit."""
        pass

    def __repr__(self) -> str:
        template = self.template_path.name if self.template_path else "blank"
        return f"PowerPresentation(template={template}, slides={self.slide_count})"
