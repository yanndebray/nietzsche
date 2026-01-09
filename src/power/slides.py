"""Slide builder utilities for Power."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from power.styles import BRAND, FONTS, POSITIONS


class SlideBuilder:
    """Builder class for adding content to a slide.

    Provides a fluent interface for adding shapes, text, tables,
    charts, and images to a PowerPoint slide.

    Example:
        >>> slide_builder = ppt.add_slide(5)
        >>> slide_builder.set_title("Sales Report")
        >>> slide_builder.add_text_box("Q1 Results", Inches(1), Inches(2))
        >>> slide_builder.add_bar_chart(categories, series_data)
    """

    def __init__(self, slide: Slide):
        """Initialize with a slide object.

        Args:
            slide: A python-pptx Slide object.
        """
        self._slide = slide

    @property
    def slide(self) -> Slide:
        """Access the underlying slide object."""
        return self._slide

    def set_title(self, text: str) -> SlideBuilder:
        """Set the slide title using the title placeholder.

        Args:
            text: Title text.

        Returns:
            Self for method chaining.
        """
        if self._slide.shapes.title:
            self._slide.shapes.title.text = text
        else:
            # Add title as text box if no placeholder
            self.add_text_box(
                text,
                POSITIONS.TITLE_LEFT,
                POSITIONS.TITLE_TOP,
                POSITIONS.TITLE_WIDTH,
                POSITIONS.TITLE_HEIGHT,
                font_size=FONTS.SLIDE_TITLE,
                bold=True,
            )
        return self

    def set_subtitle(self, text: str) -> SlideBuilder:
        """Set subtitle text (placeholder index 1).

        Args:
            text: Subtitle text.

        Returns:
            Self for method chaining.
        """
        if len(self._slide.placeholders) > 1:
            self._slide.placeholders[1].text = text
        return self

    def add_text_box(
        self,
        text: str,
        left: Inches | int = POSITIONS.CONTENT_LEFT,
        top: Inches | int = POSITIONS.CONTENT_TOP,
        width: Inches | int = POSITIONS.CONTENT_WIDTH,
        height: Inches | int = Inches(1),
        font_size: Pt | int = FONTS.BODY,
        font_name: str = "Calibri",
        bold: bool = False,
        italic: bool = False,
        color: RGBColor = BRAND.SECONDARY,
        alignment: PP_ALIGN = PP_ALIGN.LEFT,
        vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
        word_wrap: bool = True,
    ) -> SlideBuilder:
        """Add a text box to the slide.

        Args:
            text: Text content.
            left: Left position.
            top: Top position.
            width: Box width.
            height: Box height.
            font_size: Font size in points.
            font_name: Font family name.
            bold: Whether text is bold.
            italic: Whether text is italic.
            color: Text color as RGBColor.
            alignment: Text alignment.
            vertical_anchor: Vertical text anchor.
            word_wrap: Whether to wrap text.

        Returns:
            Self for method chaining.
        """
        text_box = self._slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        tf.word_wrap = word_wrap

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.name = font_name
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.alignment = alignment
        tf.vertical_anchor = vertical_anchor

        return self

    def add_bullets(
        self,
        items: list[str],
        left: Inches | int = POSITIONS.CONTENT_LEFT,
        top: Inches | int = POSITIONS.CONTENT_TOP,
        width: Inches | int = POSITIONS.CONTENT_WIDTH,
        height: Inches | int = POSITIONS.CONTENT_HEIGHT,
        font_size: Pt | int = FONTS.BODY,
        color: RGBColor = BRAND.SECONDARY,
        levels: list[int] | None = None,
    ) -> SlideBuilder:
        """Add bullet points to the slide.

        Args:
            items: List of bullet point texts.
            left: Left position.
            top: Top position.
            width: Box width.
            height: Box height.
            font_size: Font size.
            color: Text color.
            levels: Optional list of indent levels (0-based) for each item.

        Returns:
            Self for method chaining.
        """
        # Try to use content placeholder first
        content_placeholder = None
        for placeholder in self._slide.placeholders:
            if placeholder.placeholder_format.idx == 1:
                content_placeholder = placeholder
                break

        if content_placeholder and content_placeholder.has_text_frame:
            tf = content_placeholder.text_frame
            tf.clear()
        else:
            text_box = self._slide.shapes.add_textbox(left, top, width, height)
            tf = text_box.text_frame

        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = item
            p.font.size = font_size
            p.font.color.rgb = color

            if levels and i < len(levels):
                p.level = levels[i]

        return self

    def add_table(
        self,
        data: list[list[Any]],
        headers: list[str] | None = None,
        left: Inches | int = POSITIONS.TABLE_LEFT,
        top: Inches | int = POSITIONS.TABLE_TOP,
        width: Inches | int = POSITIONS.TABLE_WIDTH,
        height: Inches | int | None = None,
        header_color: RGBColor = BRAND.PRIMARY,
        header_text_color: RGBColor = BRAND.WHITE,
        font_size: Pt | int = FONTS.CAPTION,
    ) -> SlideBuilder:
        """Add a data table to the slide.

        Args:
            data: 2D list of cell values.
            headers: Optional list of header strings.
            left: Left position.
            top: Top position.
            width: Table width.
            height: Table height (auto-calculated if None).
            header_color: Background color for header row.
            header_text_color: Text color for header row.
            font_size: Font size for table cells.

        Returns:
            Self for method chaining.
        """
        rows = len(data) + (1 if headers else 0)
        cols = len(headers) if headers else (len(data[0]) if data else 1)

        if height is None:
            height = Inches(0.4 * rows)

        table_shape = self._slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        row_offset = 0
        if headers:
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = str(header)
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color

                para = cell.text_frame.paragraphs[0]
                para.font.bold = True
                para.font.color.rgb = header_text_color
                para.font.size = font_size
            row_offset = 1

        for row_idx, row_data in enumerate(data):
            for col_idx, value in enumerate(row_data):
                if col_idx >= cols:
                    break
                cell = table.cell(row_idx + row_offset, col_idx)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = font_size

        return self

    def add_bar_chart(
        self,
        categories: list[str],
        series_data: dict[str, list[float]],
        left: Inches | int = POSITIONS.CHART_LEFT,
        top: Inches | int = POSITIONS.CHART_TOP,
        width: Inches | int = POSITIONS.CHART_WIDTH,
        height: Inches | int = POSITIONS.CHART_HEIGHT,
        chart_type: XL_CHART_TYPE = XL_CHART_TYPE.COLUMN_CLUSTERED,
        has_legend: bool = True,
    ) -> SlideBuilder:
        """Add a bar/column chart to the slide.

        Args:
            categories: List of category labels.
            series_data: Dict mapping series names to value lists.
            left: Left position.
            top: Top position.
            width: Chart width.
            height: Chart height.
            chart_type: Type of chart (default: clustered column).
            has_legend: Whether to show legend.

        Returns:
            Self for method chaining.
        """
        chart_data = CategoryChartData()
        chart_data.categories = categories

        for series_name, values in series_data.items():
            if len(values) != len(categories):
                raise ValueError(
                    f"Series '{series_name}' has {len(values)} values, "
                    f"expected {len(categories)}"
                )
            chart_data.add_series(series_name, values)

        graphic_frame = self._slide.shapes.add_chart(
            chart_type, left, top, width, height, chart_data
        )

        chart = graphic_frame.chart
        chart.has_legend = has_legend
        if has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False

        return self

    def add_line_chart(
        self,
        categories: list[str],
        series_data: dict[str, list[float]],
        left: Inches | int = POSITIONS.CHART_LEFT,
        top: Inches | int = POSITIONS.CHART_TOP,
        width: Inches | int = POSITIONS.CHART_WIDTH,
        height: Inches | int = POSITIONS.CHART_HEIGHT,
    ) -> SlideBuilder:
        """Add a line chart to the slide.

        Args:
            categories: List of category labels.
            series_data: Dict mapping series names to value lists.
            left: Left position.
            top: Top position.
            width: Chart width.
            height: Chart height.

        Returns:
            Self for method chaining.
        """
        return self.add_bar_chart(
            categories,
            series_data,
            left,
            top,
            width,
            height,
            chart_type=XL_CHART_TYPE.LINE,
        )

    def add_pie_chart(
        self,
        categories: list[str],
        values: list[float],
        left: Inches | int = POSITIONS.CHART_LEFT,
        top: Inches | int = POSITIONS.CHART_TOP,
        width: Inches | int = POSITIONS.CHART_WIDTH,
        height: Inches | int = POSITIONS.CHART_HEIGHT,
    ) -> SlideBuilder:
        """Add a pie chart to the slide.

        Args:
            categories: List of category labels.
            values: List of values for each category.
            left: Left position.
            top: Top position.
            width: Chart width.
            height: Chart height.

        Returns:
            Self for method chaining.
        """
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series("Series 1", values)

        graphic_frame = self._slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, left, top, width, height, chart_data
        )

        chart = graphic_frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT

        return self

    def add_image(
        self,
        image_path: str | Path | BytesIO,
        left: Inches | int = POSITIONS.CONTENT_LEFT,
        top: Inches | int = POSITIONS.CONTENT_TOP,
        width: Inches | int | None = None,
        height: Inches | int | None = None,
    ) -> SlideBuilder:
        """Add an image to the slide.

        Args:
            image_path: Path to image file or BytesIO object.
            left: Left position.
            top: Top position.
            width: Image width (maintains aspect ratio if height not given).
            height: Image height (maintains aspect ratio if width not given).

        Returns:
            Self for method chaining.
        """
        if isinstance(image_path, (str, Path)):
            image_path = Path(image_path)
            if not image_path.exists():
                raise FileNotFoundError(f"Image not found: {image_path}")
            image_source = str(image_path)
        else:
            image_source = image_path

        # Default width if neither dimension specified
        if width is None and height is None:
            width = Inches(4)

        self._slide.shapes.add_picture(image_source, left, top, width=width, height=height)
        return self

    def add_shape(
        self,
        shape_type: MSO_SHAPE = MSO_SHAPE.RECTANGLE,
        left: Inches | int = Inches(1),
        top: Inches | int = Inches(1),
        width: Inches | int = Inches(2),
        height: Inches | int = Inches(1),
        fill_color: RGBColor | None = BRAND.PRIMARY,
        line_color: RGBColor | None = None,
        text: str | None = None,
        font_color: RGBColor = BRAND.WHITE,
        font_size: Pt | int = FONTS.BODY,
    ) -> SlideBuilder:
        """Add a shape to the slide.

        Args:
            shape_type: Type of shape (e.g., MSO_SHAPE.RECTANGLE).
            left: Left position.
            top: Top position.
            width: Shape width.
            height: Shape height.
            fill_color: Fill color (None for no fill).
            line_color: Line/border color (None for no line).
            text: Optional text inside the shape.
            font_color: Text color if text is provided.
            font_size: Font size if text is provided.

        Returns:
            Self for method chaining.
        """
        shape = self._slide.shapes.add_shape(shape_type, left, top, width, height)

        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()

        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background()

        if text:
            shape.text = text
            para = shape.text_frame.paragraphs[0]
            para.font.color.rgb = font_color
            para.font.size = font_size
            para.alignment = PP_ALIGN.CENTER
            shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        return self

    def add_note(self, text: str) -> SlideBuilder:
        """Add speaker notes to the slide.

        Args:
            text: Notes text.

        Returns:
            Self for method chaining.
        """
        notes_slide = self._slide.notes_slide
        notes_slide.notes_text_frame.text = text
        return self
