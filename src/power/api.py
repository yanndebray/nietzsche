"""FastAPI server for Power presentation generator.

Run locally:
    uvicorn power.api:app --reload

Deploy to Cloud Run:
    See Dockerfile and deploy.sh
"""

from __future__ import annotations

import io
import logging
import tempfile
from pathlib import Path
from typing import Any

from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field

from power.core import PowerPresentation

logger = logging.getLogger(__name__)

app = FastAPI(
    title="Power API",
    description="Generate PowerPoint presentations from YAML/JSON specifications",
    version="0.1.7",
    docs_url="/",
)


# --- Pydantic Models ---

class SlideSpec(BaseModel):
    """Specification for a single slide."""
    type: str = Field(default="content", description="Slide type: title, section, content, table, chart, image, blank")
    title: str | None = None
    subtitle: str | None = None
    bullets: list[str] | None = None
    layout: str | int | None = None
    headers: list[str] | None = None
    data: list[list[Any]] | None = None
    chart_type: str | None = Field(default=None, description="Chart type: bar, line, pie")
    categories: list[str] | None = None
    series: dict[str, list[float]] | None = None
    image: str | None = None
    image_left: float | None = None
    image_top: float | None = None
    image_width: float | None = None
    image_height: float | None = None
    left: float | None = None
    top: float | None = None
    width: float | None = None
    height: float | None = None
    note: str | None = None
    placeholder_images: dict[str, str] | None = None


class PresentationSpec(BaseModel):
    """Full presentation specification."""
    title: str | None = Field(default=None, description="Title slide title")
    subtitle: str | None = Field(default=None, description="Title slide subtitle")
    slides: list[SlideSpec] = Field(default_factory=list)

    model_config = {
        "json_schema_extra": {
            "examples": [
                {
                    "title": "My Presentation",
                    "subtitle": "Created with Power API",
                    "slides": [
                        {
                            "type": "content",
                            "title": "Key Points",
                            "bullets": ["First point", "Second point", "Third point"]
                        },
                        {
                            "type": "section",
                            "title": "Next Section"
                        },
                        {
                            "type": "chart",
                            "title": "Sales Data",
                            "chart_type": "bar",
                            "categories": ["Q1", "Q2", "Q3", "Q4"],
                            "series": {"Revenue": [100, 150, 120, 180]}
                        }
                    ]
                }
            ]
        }
    }


class HealthResponse(BaseModel):
    """Health check response."""
    status: str
    version: str


# --- Helper Functions ---

def _build_from_spec(ppt: PowerPresentation, spec: dict[str, Any]) -> None:
    """Build presentation from specification dict.

    This is adapted from cli.py to work with dict input.
    """
    from pptx.util import Inches

    # Handle title slide if top-level title exists
    if spec.get("title"):
        ppt.add_title_slide(spec.get("title", ""), spec.get("subtitle", ""))

    # Process each slide
    for slide_spec in spec.get("slides", []):
        slide_type = slide_spec.get("type", "content")
        builder = None

        if slide_type == "title":
            builder = ppt.add_title_slide(
                slide_spec.get("title", ""),
                slide_spec.get("subtitle", ""),
            )

        elif slide_type == "section":
            builder = ppt.add_section_slide(
                slide_spec.get("title", ""),
                slide_spec.get("subtitle", ""),
            )

        elif slide_type == "content":
            layout = slide_spec.get("layout")
            builder = ppt.add_content_slide(
                slide_spec.get("title", ""),
                slide_spec.get("bullets"),
                layout=layout,
            )
            # Support optional image in content slides
            image_path = slide_spec.get("image")
            if image_path:
                left = slide_spec.get("image_left", 7)
                top = slide_spec.get("image_top", 1.5)
                width = slide_spec.get("image_width", 5)
                height = slide_spec.get("image_height")

                kwargs = {
                    "left": Inches(left),
                    "top": Inches(top),
                    "width": Inches(width),
                }
                if height is not None:
                    kwargs["height"] = Inches(height)

                builder.add_image(image_path, **kwargs)

        elif slide_type == "table":
            title_only_layout = ppt.find_title_only_layout()
            builder = ppt.add_slide(title_only_layout)
            builder.set_title(slide_spec.get("title", ""))
            builder.add_table(
                slide_spec.get("data", []),
                slide_spec.get("headers"),
            )

        elif slide_type == "chart":
            title_only_layout = ppt.find_title_only_layout()
            builder = ppt.add_slide(title_only_layout)
            builder.set_title(slide_spec.get("title", ""))
            chart_type = slide_spec.get("chart_type", "bar")
            categories = slide_spec.get("categories", [])
            series = slide_spec.get("series", {})

            if chart_type == "pie":
                values = list(series.values())[0] if series else []
                builder.add_pie_chart(categories, values)
            elif chart_type == "line":
                builder.add_line_chart(categories, series)
            else:
                builder.add_bar_chart(categories, series)

        elif slide_type == "image":
            title_only_layout = ppt.find_title_only_layout()
            builder = ppt.add_slide(title_only_layout)
            if slide_spec.get("title"):
                builder.set_title(slide_spec["title"])

            image_path = slide_spec.get("image") or slide_spec.get("path")
            if image_path:
                left = slide_spec.get("left")
                top = slide_spec.get("top")
                width = slide_spec.get("width")
                height = slide_spec.get("height")

                left = Inches(left) if left is not None else None
                top = Inches(top) if top is not None else None
                width = Inches(width) if width is not None else None
                height = Inches(height) if height is not None else None

                kwargs = {}
                if left is not None:
                    kwargs["left"] = left
                if top is not None:
                    kwargs["top"] = top
                if width is not None:
                    kwargs["width"] = width
                if height is not None:
                    kwargs["height"] = height

                builder.add_image(image_path, **kwargs)

        elif slide_type == "blank":
            builder = ppt.add_blank_slide()

        # Add speaker notes if specified
        if builder and slide_spec.get("note"):
            builder.add_note(slide_spec["note"])


def _generate_pptx_bytes(spec: dict[str, Any], template_bytes: bytes | None = None) -> bytes:
    """Generate presentation and return as bytes."""
    with tempfile.TemporaryDirectory() as tmpdir:
        template_path = None

        # Save template if provided
        if template_bytes:
            template_path = Path(tmpdir) / "template.pptx"
            template_path.write_bytes(template_bytes)

        # Create presentation
        ppt = PowerPresentation(template=template_path)

        if template_path:
            ppt.clear_slides()

        # Build from spec
        _build_from_spec(ppt, spec)

        # Save to bytes
        output_path = Path(tmpdir) / "output.pptx"
        ppt.save(output_path)

        return output_path.read_bytes()


# --- API Endpoints ---

@app.get("/health", response_model=HealthResponse, tags=["System"])
async def health_check():
    """Health check endpoint for Cloud Run."""
    return HealthResponse(status="healthy", version="0.1.7")


@app.post("/generate", tags=["Generate"])
async def generate_presentation(
    spec: PresentationSpec,
    filename: str = "presentation.pptx"
) -> StreamingResponse:
    """Generate a PowerPoint presentation from a JSON specification.

    Returns the generated .pptx file as a download.

    Example request body:
    ```json
    {
        "title": "My Presentation",
        "subtitle": "Created with Power API",
        "slides": [
            {
                "type": "content",
                "title": "Key Points",
                "bullets": ["First point", "Second point"]
            }
        ]
    }
    ```
    """
    try:
        spec_dict = spec.model_dump(exclude_none=True)
        pptx_bytes = _generate_pptx_bytes(spec_dict)

        return StreamingResponse(
            io.BytesIO(pptx_bytes),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f'attachment; filename="{filename}"'
            }
        )
    except Exception as e:
        logger.exception("Failed to generate presentation")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/generate-with-template", tags=["Generate"])
async def generate_with_template(
    template: UploadFile = File(..., description="Template .pptx file"),
    spec: str = Form(..., description="JSON specification for slides"),
    filename: str = Form(default="presentation.pptx", description="Output filename")
) -> StreamingResponse:
    """Generate a presentation using an uploaded template.

    Upload a .pptx template file and provide a JSON specification.
    The template's styling will be preserved.
    """
    import json

    try:
        # Validate template file
        if not template.filename or not template.filename.endswith(".pptx"):
            raise HTTPException(status_code=400, detail="Template must be a .pptx file")

        # Parse spec
        try:
            spec_dict = json.loads(spec)
        except json.JSONDecodeError as e:
            raise HTTPException(status_code=400, detail=f"Invalid JSON spec: {e}")

        # Read template
        template_bytes = await template.read()

        # Generate
        pptx_bytes = _generate_pptx_bytes(spec_dict, template_bytes)

        return StreamingResponse(
            io.BytesIO(pptx_bytes),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f'attachment; filename="{filename}"'
            }
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("Failed to generate presentation with template")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/inspect", tags=["Utilities"])
async def inspect_template(
    template: UploadFile = File(..., description="Template .pptx file to inspect")
) -> dict[str, Any]:
    """Inspect a PowerPoint template.

    Upload a .pptx file to see its available layouts and placeholders.
    Useful for understanding how to use the template in generation.
    """
    try:
        if not template.filename or not template.filename.endswith(".pptx"):
            raise HTTPException(status_code=400, detail="File must be a .pptx")

        with tempfile.TemporaryDirectory() as tmpdir:
            template_path = Path(tmpdir) / "template.pptx"
            template_path.write_bytes(await template.read())

            ppt = PowerPresentation(template=template_path)
            info = ppt.get_template_info()

            # Convert non-JSON-serializable values
            info["slide_width"] = str(info["slide_width"])
            info["slide_height"] = str(info["slide_height"])
            info["template"] = template.filename

            return info
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("Failed to inspect template")
        raise HTTPException(status_code=500, detail=str(e))


# --- Batch Processing ---

class BatchItem(BaseModel):
    """Single item in a batch generation request."""
    spec: PresentationSpec
    filename: str = "presentation.pptx"


class BatchRequest(BaseModel):
    """Batch generation request."""
    items: list[BatchItem] = Field(..., max_length=50, description="Up to 50 presentations per batch")


@app.post("/batch/generate", tags=["Batch"])
async def batch_generate(request: BatchRequest) -> dict[str, Any]:
    """Generate multiple presentations in a single request.

    Returns a summary with download URLs (for Cloud Storage integration)
    or base64-encoded files for smaller batches.

    Note: For large batches, consider using Cloud Tasks or Pub/Sub
    for async processing.
    """
    import base64

    results = []
    errors = []

    for i, item in enumerate(request.items):
        try:
            spec_dict = item.spec.model_dump(exclude_none=True)
            pptx_bytes = _generate_pptx_bytes(spec_dict)

            results.append({
                "index": i,
                "filename": item.filename,
                "size_bytes": len(pptx_bytes),
                "data_base64": base64.b64encode(pptx_bytes).decode("utf-8")
            })
        except Exception as e:
            errors.append({
                "index": i,
                "filename": item.filename,
                "error": str(e)
            })

    return {
        "total": len(request.items),
        "successful": len(results),
        "failed": len(errors),
        "results": results,
        "errors": errors
    }


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8080)
