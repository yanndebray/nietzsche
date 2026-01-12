"""CLI interface for Power presentation generator."""

from __future__ import annotations

import json
import sys
from pathlib import Path
from typing import Any

import click
import yaml
from rich.console import Console
from rich.panel import Panel
from rich.table import Table
from rich.tree import Tree

from power.core import PowerPresentation
from power.styles import BRAND, FONTS

console = Console()


@click.group(context_settings={"help_option_names": ["-h", "--help"]})
@click.version_option(prog_name="power", package_name="nietzsche")
def main():
    """Power - Generate PowerPoint presentations from the command line.

    A CLI tool for creating and manipulating PowerPoint slide decks
    programmatically. Supports templates, YAML/JSON input, and various
    slide types including titles, content, tables, and charts.

    Examples:

        # Create a simple presentation
        power create output.pptx --title "My Presentation"

        # Use a template
        power create output.pptx --template Galaxy.pptx --title "Themed Deck"

        # Generate from YAML/JSON file
        power generate slides.yaml -o presentation.pptx

        # Inspect a template
        power inspect Galaxy.pptx

        # Create presentation interactively
        power new
    """
    pass


@main.command()
@click.argument("output", type=click.Path())
@click.option("-t", "--template", type=click.Path(exists=True), help="Template .pptx file")
@click.option("--title", default="Untitled Presentation", help="Presentation title")
@click.option("--subtitle", default="", help="Presentation subtitle")
@click.option("--author", default="", help="Author name")
def create(output: str, template: str | None, title: str, subtitle: str, author: str):
    """Create a new presentation with a title slide.

    OUTPUT is the path for the generated .pptx file.

    Examples:

        power create my_deck.pptx --title "Q1 Report"

        power create themed.pptx -t Galaxy.pptx --title "Galaxy Theme"
    """
    try:
        ppt = PowerPresentation(template=template)

        # If using template, clear existing slides
        if template:
            ppt.clear_slides()

        # Add title slide
        ppt.add_title_slide(title, subtitle)

        # Add author as note if provided
        if author:
            ppt.presentation.core_properties.author = author

        ppt.save(output)
        console.print(f"[green]Created presentation:[/green] {output}")

    except Exception as e:
        console.print(f"[red]Error:[/red] {e}")
        sys.exit(1)


@main.command()
@click.argument("input_file", type=click.Path(exists=True))
@click.option("-o", "--output", required=True, type=click.Path(), help="Output .pptx file")
@click.option("-t", "--template", type=click.Path(exists=True), help="Template .pptx file")
def generate(input_file: str, output: str, template: str | None):
    """Generate a presentation from a YAML or JSON file.

    INPUT_FILE is a YAML or JSON file describing the slides.

    Supported slide types: title, section, content, table, chart, image, blank

    Example YAML format:

    \b
        title: My Presentation
        subtitle: A Great Deck
        slides:
          - type: content
            title: Key Points
            layout: "Title and Content 1"  # optional layout name
            bullets:
              - First point
              - Second point
            image: path/to/image.png       # optional free-positioned image
            placeholder_images:            # fill template placeholders
              10: path/to/image.png        # placeholder idx: image path
          - type: image
            title: Full Image Slide
            image: path/to/image.png
            width: 8                       # optional (inches)
          - type: table
            title: Data Table
            headers: [Name, Value]
            data:
              - [Item A, 100]

    Use 'power inspect template.pptx' to see available layouts and placeholder indices.
    """
    try:
        # Load input file
        input_path = Path(input_file)
        with open(input_path) as f:
            if input_path.suffix.lower() in (".yaml", ".yml"):
                spec = yaml.safe_load(f)
            else:
                spec = json.load(f)

        ppt = PowerPresentation(template=template)

        if template:
            ppt.clear_slides()

        # Process slides
        _build_from_spec(ppt, spec)

        ppt.save(output)
        console.print(f"[green]Generated presentation:[/green] {output}")
        console.print(f"  Slides: {ppt.slide_count}")

    except Exception as e:
        console.print(f"[red]Error:[/red] {e}")
        sys.exit(1)


def _fill_placeholder_images(builder, slide_spec: dict[str, Any]) -> None:
    """Fill picture placeholders from slide spec."""
    placeholder_images = slide_spec.get("placeholder_images", {})
    if not placeholder_images:
        return

    for idx_str, image_path in placeholder_images.items():
        try:
            idx = int(idx_str)
            builder.fill_picture_placeholder(image_path, placeholder_idx=idx)
        except (ValueError, FileNotFoundError) as e:
            # Log warning but continue with other placeholders
            console.print(f"[yellow]Warning:[/yellow] Could not fill placeholder {idx_str}: {e}")


def _build_from_spec(ppt: PowerPresentation, spec: dict[str, Any]) -> None:
    """Build presentation from specification dict."""
    # Handle title slide if top-level title exists
    if "title" in spec:
        ppt.add_title_slide(spec.get("title", ""), spec.get("subtitle", ""))

    # Process each slide
    for slide_spec in spec.get("slides", []):
        slide_type = slide_spec.get("type", "content")
        builder = None

        if slide_type == "title":
            builder = ppt.add_title_slide(
                slide_spec.get("title", ""),
                slide_spec.get("subtitle", ""),
            )

        elif slide_type == "section":
            builder = ppt.add_section_slide(
                slide_spec.get("title", ""),
                slide_spec.get("subtitle", ""),
            )

        elif slide_type == "content":
            # Support optional layout specification
            layout = slide_spec.get("layout")
            builder = ppt.add_content_slide(
                slide_spec.get("title", ""),
                slide_spec.get("bullets"),
                layout=layout,
            )
            # Support optional image in content slides
            image_path = slide_spec.get("image")
            if image_path:
                from pptx.util import Inches

                # Position image on the right side by default for content slides
                left = slide_spec.get("image_left", 7)
                top = slide_spec.get("image_top", 1.5)
                width = slide_spec.get("image_width", 5)
                height = slide_spec.get("image_height")

                kwargs = {
                    "left": Inches(left),
                    "top": Inches(top),
                    "width": Inches(width),
                }
                if height is not None:
                    kwargs["height"] = Inches(height)

                builder.add_image(image_path, **kwargs)

            if "note" in slide_spec:
                builder.add_note(slide_spec["note"])

        elif slide_type == "table":
            title_only_layout = ppt.find_title_only_layout()
            builder = ppt.add_slide(title_only_layout)
            builder.set_title(slide_spec.get("title", ""))
            builder.add_table(
                slide_spec.get("data", []),
                slide_spec.get("headers"),
            )

        elif slide_type == "chart":
            title_only_layout = ppt.find_title_only_layout()
            builder = ppt.add_slide(title_only_layout)
            builder.set_title(slide_spec.get("title", ""))
            chart_type = slide_spec.get("chart_type", "bar")
            categories = slide_spec.get("categories", [])
            series = slide_spec.get("series", {})

            if chart_type == "pie":
                values = list(series.values())[0] if series else []
                builder.add_pie_chart(categories, values)
            elif chart_type == "line":
                builder.add_line_chart(categories, series)
            else:
                builder.add_bar_chart(categories, series)

        elif slide_type == "image":
            title_only_layout = ppt.find_title_only_layout()
            builder = ppt.add_slide(title_only_layout)
            if "title" in slide_spec:
                builder.set_title(slide_spec["title"])

            image_path = slide_spec.get("image") or slide_spec.get("path")
            if image_path:
                # Convert position/size values to Inches if provided
                from pptx.util import Inches

                left = slide_spec.get("left")
                top = slide_spec.get("top")
                width = slide_spec.get("width")
                height = slide_spec.get("height")

                # Convert numeric values to Inches
                left = Inches(left) if left is not None else None
                top = Inches(top) if top is not None else None
                width = Inches(width) if width is not None else None
                height = Inches(height) if height is not None else None

                # Build kwargs, only including non-None values
                kwargs = {}
                if left is not None:
                    kwargs["left"] = left
                if top is not None:
                    kwargs["top"] = top
                if width is not None:
                    kwargs["width"] = width
                if height is not None:
                    kwargs["height"] = height

                builder.add_image(image_path, **kwargs)

            if "note" in slide_spec:
                builder.add_note(slide_spec["note"])

        elif slide_type == "blank":
            builder = ppt.add_blank_slide()

        # Fill picture placeholders if specified (works with any slide type)
        if builder and "placeholder_images" in slide_spec:
            _fill_placeholder_images(builder, slide_spec)


@main.command()
@click.argument("pptx_file", type=click.Path(exists=True))
@click.option("--json", "as_json", is_flag=True, help="Output as JSON")
def inspect(pptx_file: str, as_json: bool):
    """Inspect a PowerPoint template or presentation.

    Shows available layouts, placeholders, and slide information.

    Examples:

        power inspect Galaxy.pptx

        power inspect template.pptx --json
    """
    try:
        ppt = PowerPresentation(template=pptx_file)
        info = ppt.get_template_info()

        if as_json:
            # Convert non-serializable values
            info["slide_width"] = str(info["slide_width"])
            info["slide_height"] = str(info["slide_height"])
            click.echo(json.dumps(info, indent=2))
            return

        # Rich formatted output
        console.print()
        console.print(Panel(f"[bold]{pptx_file}[/bold]", title="PowerPoint Inspection"))

        # Basic info
        info_table = Table(show_header=False, box=None)
        info_table.add_row("Slides:", str(info["slide_count"]))
        info_table.add_row("Layouts:", str(info["layout_count"]))
        info_table.add_row(
            "Dimensions:",
            f"{info['slide_width'].inches:.2f}\" x {info['slide_height'].inches:.2f}\""
        )
        console.print(info_table)
        console.print()

        # Layout tree
        tree = Tree("[bold]Available Layouts[/bold]")
        for layout in info["layouts"]:
            layout_branch = tree.add(f"[cyan]{layout['index']}[/cyan]: {layout['name']}")
            if layout["placeholders"]:
                for ph in layout["placeholders"]:
                    ph_type = ph["type"].replace("PLACEHOLDER_FORMAT_TYPE.", "")
                    layout_branch.add(f"[dim]idx={ph['idx']}[/dim] {ph['name']} ({ph_type})")

        console.print(tree)

    except Exception as e:
        console.print(f"[red]Error:[/red] {e}")
        sys.exit(1)


@main.command()
@click.option("-t", "--template", type=click.Path(exists=True), help="Template .pptx file")
@click.option("-o", "--output", default="presentation.pptx", help="Output file path")
def new(template: str | None, output: str):
    """Interactively create a new presentation.

    Guides you through creating slides one at a time.

    Examples:

        power new

        power new -t Galaxy.pptx -o my_deck.pptx
    """
    try:
        console.print(Panel("[bold]Power - Interactive Presentation Builder[/bold]"))

        ppt = PowerPresentation(template=template)
        if template:
            ppt.clear_slides()
            console.print(f"Using template: [cyan]{template}[/cyan]")

        # Get title
        title = click.prompt("Presentation title", default="Untitled")
        subtitle = click.prompt("Subtitle (optional)", default="")
        ppt.add_title_slide(title, subtitle)

        while True:
            console.print("\n[bold]Add a slide:[/bold]")
            console.print("  1. Content slide (bullets)")
            console.print("  2. Section header")
            console.print("  3. Blank slide")
            console.print("  4. Done - save presentation")

            choice = click.prompt("Choice", type=int, default=4)

            if choice == 1:
                slide_title = click.prompt("Slide title")
                console.print("Enter bullet points (empty line to finish):")
                bullets = []
                while True:
                    bullet = click.prompt("  -", default="")
                    if not bullet:
                        break
                    bullets.append(bullet)
                ppt.add_content_slide(slide_title, bullets if bullets else None)
                console.print("[green]Added content slide[/green]")

            elif choice == 2:
                section_title = click.prompt("Section title")
                section_subtitle = click.prompt("Subtitle (optional)", default="")
                ppt.add_section_slide(section_title, section_subtitle)
                console.print("[green]Added section slide[/green]")

            elif choice == 3:
                ppt.add_blank_slide()
                console.print("[green]Added blank slide[/green]")

            elif choice == 4:
                break

        ppt.save(output)
        console.print(f"\n[green]Saved:[/green] {output} ({ppt.slide_count} slides)")

    except click.Abort:
        console.print("\n[yellow]Cancelled[/yellow]")
        sys.exit(0)
    except Exception as e:
        console.print(f"[red]Error:[/red] {e}")
        sys.exit(1)


@main.command()
@click.argument("pptx_file", type=click.Path(exists=True))
@click.argument("replacements", nargs=-1)
@click.option("-o", "--output", required=True, help="Output file path")
def replace(pptx_file: str, replacements: tuple[str, ...], output: str):
    """Replace placeholders in a presentation.

    PPTX_FILE is the input presentation.
    REPLACEMENTS are key=value pairs for {{KEY}} placeholders.

    Examples:

        power replace template.pptx NAME=John DATE=2024-01-01 -o filled.pptx

        power replace report.pptx "TITLE=Q1 Report" -o q1.pptx
    """
    try:
        # Parse replacements
        replacement_dict = {}
        for r in replacements:
            if "=" not in r:
                raise click.BadParameter(f"Invalid replacement format: {r} (expected KEY=VALUE)")
            key, value = r.split("=", 1)
            replacement_dict[f"{{{{{key}}}}}"] = value

        ppt = PowerPresentation(template=pptx_file)
        ppt.replace_placeholders(replacement_dict)
        ppt.save(output)

        console.print(f"[green]Replaced {len(replacement_dict)} placeholder(s)[/green]")
        console.print(f"Saved to: {output}")

    except Exception as e:
        console.print(f"[red]Error:[/red] {e}")
        sys.exit(1)


@main.command()
@click.argument("pptx_file", type=click.Path(exists=True))
@click.argument("slide_index", type=int)
@click.option("-o", "--output", required=True, help="Output file path")
def remove(pptx_file: str, slide_index: int, output: str):
    """Remove a slide from a presentation.

    SLIDE_INDEX is the 0-based index of the slide to remove.

    Examples:

        power remove deck.pptx 2 -o deck_trimmed.pptx
    """
    try:
        ppt = PowerPresentation(template=pptx_file)
        original_count = ppt.slide_count

        ppt.remove_slide(slide_index)
        ppt.save(output)

        console.print(f"[green]Removed slide {slide_index}[/green]")
        console.print(f"Slides: {original_count} -> {ppt.slide_count}")

    except IndexError:
        console.print(f"[red]Error:[/red] Slide index {slide_index} out of range")
        sys.exit(1)
    except Exception as e:
        console.print(f"[red]Error:[/red] {e}")
        sys.exit(1)


@main.command(name="add")
@click.argument("pptx_file", type=click.Path(exists=True))
@click.option("-o", "--output", required=True, help="Output file path")
@click.option("--title", required=True, help="Slide title")
@click.option("--bullets", multiple=True, help="Bullet points (can specify multiple)")
@click.option("--layout", default=None, type=int, help="Layout index (auto-detect if not specified)")
def add_slide(pptx_file: str, output: str, title: str, bullets: tuple, layout: int | None):
    """Add a slide to an existing presentation.

    Examples:

        power add deck.pptx -o updated.pptx --title "New Slide" --bullets "Point 1" --bullets "Point 2"
    """
    try:
        ppt = PowerPresentation(template=pptx_file)
        ppt.add_content_slide(title, list(bullets) if bullets else None, layout=layout)
        ppt.save(output)

        console.print(f"[green]Added slide:[/green] {title}")
        console.print(f"Total slides: {ppt.slide_count}")

    except Exception as e:
        console.print(f"[red]Error:[/red] {e}")
        sys.exit(1)


# --- Cloud API Commands ---

DEFAULT_API_URL = "https://power-api-944767079044.us-central1.run.app"


@main.group()
@click.option(
    "--api-url",
    default=DEFAULT_API_URL,
    envvar="POWER_API_URL",
    help="Cloud API URL (default: Cloud Run service)",
)
@click.pass_context
def cloud(ctx, api_url: str):
    """Generate presentations via cloud API.

    Uses the deployed Power API service to generate presentations
    in the cloud instead of locally.

    Examples:

        # Check API health
        power cloud health

        # Generate via cloud
        power cloud generate slides.yaml -o presentation.pptx

        # Use custom API URL
        power cloud --api-url http://localhost:8080 health

        # Or via environment variable
        POWER_API_URL=http://localhost:8080 power cloud health
    """
    ctx.ensure_object(dict)
    ctx.obj["api_url"] = api_url


@cloud.command()
@click.pass_context
def health(ctx):
    """Check cloud API health status.

    Examples:

        power cloud health
    """
    import httpx

    api_url = ctx.obj["api_url"]

    try:
        with httpx.Client(timeout=10.0) as client:
            response = client.get(f"{api_url}/health")
            response.raise_for_status()
            data = response.json()

        console.print(f"[green]API Status:[/green] {data.get('status', 'unknown')}")
        console.print(f"[dim]Version:[/dim] {data.get('version', 'unknown')}")
        console.print(f"[dim]URL:[/dim] {api_url}")

    except httpx.ConnectError:
        console.print(f"[red]Error:[/red] Cannot connect to API at {api_url}")
        sys.exit(1)
    except httpx.TimeoutException:
        console.print(f"[red]Error:[/red] Request timed out")
        sys.exit(1)
    except httpx.HTTPStatusError as e:
        console.print(f"[red]Error:[/red] HTTP {e.response.status_code}")
        sys.exit(1)
    except Exception as e:
        console.print(f"[red]Error:[/red] {e}")
        sys.exit(1)


@cloud.command(name="generate")
@click.argument("input_file", type=click.Path(exists=True))
@click.option("-o", "--output", required=True, type=click.Path(), help="Output .pptx file")
@click.option("-t", "--template", type=click.Path(exists=True), help="Template .pptx file")
@click.pass_context
def cloud_generate(ctx, input_file: str, output: str, template: str | None):
    """Generate a presentation via cloud API.

    INPUT_FILE is a YAML or JSON file describing the slides.

    Examples:

        power cloud generate slides.yaml -o presentation.pptx

        power cloud generate slides.yaml -o themed.pptx -t Galaxy.pptx
    """
    import httpx

    api_url = ctx.obj["api_url"]

    try:
        # Load input file
        input_path = Path(input_file)
        with open(input_path) as f:
            if input_path.suffix.lower() in (".yaml", ".yml"):
                spec = yaml.safe_load(f)
            else:
                spec = json.load(f)

        with httpx.Client(timeout=60.0) as client:
            if template:
                # Use generate-with-template endpoint
                with open(template, "rb") as template_file:
                    files = {"template": (Path(template).name, template_file, "application/octet-stream")}
                    data = {
                        "spec": json.dumps(spec),
                        "filename": Path(output).name,
                    }
                    response = client.post(
                        f"{api_url}/generate-with-template",
                        files=files,
                        data=data,
                    )
            else:
                # Use generate endpoint
                response = client.post(
                    f"{api_url}/generate",
                    json=spec,
                    params={"filename": Path(output).name},
                )

            response.raise_for_status()

        # Save output
        output_path = Path(output)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_bytes(response.content)

        console.print(f"[green]Generated via cloud:[/green] {output}")
        console.print(f"[dim]Size:[/dim] {len(response.content):,} bytes")

    except httpx.ConnectError:
        console.print(f"[red]Error:[/red] Cannot connect to API at {api_url}")
        sys.exit(1)
    except httpx.TimeoutException:
        console.print(f"[red]Error:[/red] Request timed out")
        sys.exit(1)
    except httpx.HTTPStatusError as e:
        console.print(f"[red]Error:[/red] HTTP {e.response.status_code}: {e.response.text}")
        sys.exit(1)
    except Exception as e:
        console.print(f"[red]Error:[/red] {e}")
        sys.exit(1)


@cloud.command(name="inspect")
@click.argument("template", type=click.Path(exists=True))
@click.option("--json", "as_json", is_flag=True, help="Output as JSON")
@click.pass_context
def cloud_inspect(ctx, template: str, as_json: bool):
    """Inspect a template via cloud API.

    Upload a template to see its available layouts and placeholders.

    Examples:

        power cloud inspect Galaxy.pptx

        power cloud inspect template.pptx --json
    """
    import httpx

    api_url = ctx.obj["api_url"]

    try:
        with httpx.Client(timeout=30.0) as client:
            with open(template, "rb") as template_file:
                files = {"template": (Path(template).name, template_file, "application/octet-stream")}
                response = client.post(f"{api_url}/inspect", files=files)
                response.raise_for_status()

        info = response.json()

        if as_json:
            click.echo(json.dumps(info, indent=2))
            return

        # Rich formatted output
        console.print()
        console.print(Panel(f"[bold]{template}[/bold] (via cloud)", title="Template Inspection"))

        info_table = Table(show_header=False, box=None)
        info_table.add_row("Slides:", str(info.get("slide_count", "?")))
        info_table.add_row("Layouts:", str(info.get("layout_count", "?")))
        console.print(info_table)
        console.print()

        # Layout tree
        tree = Tree("[bold]Available Layouts[/bold]")
        for layout in info.get("layouts", []):
            layout_branch = tree.add(f"[cyan]{layout['index']}[/cyan]: {layout['name']}")
            for ph in layout.get("placeholders", []):
                ph_type = ph["type"].replace("PLACEHOLDER_FORMAT_TYPE.", "")
                layout_branch.add(f"[dim]idx={ph['idx']}[/dim] {ph['name']} ({ph_type})")

        console.print(tree)

    except httpx.ConnectError:
        console.print(f"[red]Error:[/red] Cannot connect to API at {api_url}")
        sys.exit(1)
    except httpx.TimeoutException:
        console.print(f"[red]Error:[/red] Request timed out")
        sys.exit(1)
    except httpx.HTTPStatusError as e:
        console.print(f"[red]Error:[/red] HTTP {e.response.status_code}: {e.response.text}")
        sys.exit(1)
    except Exception as e:
        console.print(f"[red]Error:[/red] {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
